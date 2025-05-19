import os
import requests
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
from docx import Document as DocxDocument  # Alias to avoid conflicts
from pptx import Presentation  # For PowerPoint presentations
from langchain.schema import Document

# Load environment variables from .env
load_dotenv()

# Hugging Face API URL and Token
HUGGINGFACE_API_URL = os.getenv("HUGGINGFACE_API_URL")
HUGGINGFACE_API_TOKEN = os.getenv("HUGGINGFACE_API_KEY")  # Load token from .env

# Validate the Hugging Face API token
if not HUGGINGFACE_API_TOKEN:
    raise ValueError("HUGGINGFACE_API_TOKEN is not set. Please check your .env file.")

# Function to load and extract text from a Word document
def load_word(file_path):
    doc = DocxDocument(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
    return [Document(page_content=content, metadata={"source": file_path})]

# Function to load and extract text from a PowerPoint presentation
def load_pptx(file_path):
    presentation = Presentation(file_path)
    slides_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        slides_content.append("\n".join(slide_text))
    content = "\n".join(slides_content)
    return [Document(page_content=content, metadata={"source": file_path})]

# Function to load and extract text from a text file
def load_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()
    return [Document(page_content=content, metadata={"source": file_path})]

# Function to load and extract text from a PDF
def load_pdf(file_path):
    loader = PyMuPDFLoader(file_path)
    documents = loader.load()
    return [Document(page_content=doc.page_content, metadata={"source": file_path}) for doc in documents]

# Function to split text into manageable chunks
def split_text(documents):
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_documents(documents)
    return chunks
pip 
# Function to summarize text chunks using Hugging Face API
def summarize_chunks(chunks):
    headers = {"Authorization": f"Bearer {HUGGINGFACE_API_TOKEN}"}
    summaries = []
    for chunk in chunks:
        payload = {"inputs": chunk.page_content}  # Access the page_content attribute
        response = requests.post(HUGGINGFACE_API_URL, headers=headers, json=payload)
        if response.status_code == 200:
            summary = response.json()[0]['summary_text']
            summaries.append(summary)
        else:
            print(f"Error: {response.status_code}, {response.text}")
            summaries.append("Error in summarization.")
    return summaries

# Function to process all supported files in a folder
def process_folder(folder_path):
    all_documents = []
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        print(f"Processing file: {filename}")
        if filename.endswith(".pdf"):
            documents = load_pdf(file_path)
        elif filename.endswith(".docx"):
            documents = load_word(file_path)
        elif filename.endswith(".pptx"):
            documents = load_pptx(file_path)
        elif filename.endswith(".txt"):
            documents = load_txt(file_path)
        else:
            print(f"Skipping unsupported file: {filename}")
            continue
        all_documents.extend(documents)
    return all_documents

# Main function to summarize all files in a folder
def summarize_folder(folder_path):
    print(f"Loading files from folder: {folder_path}")
    documents = process_folder(folder_path)
    print("Splitting text into chunks...")
    chunks = split_text(documents)
    print("Generating summaries...")
    summaries = summarize_chunks(chunks)
    print("\nFinal Summary:")
    print("\n".join(summaries))

if __name__ == "__main__":
    # Replace 'your_folder_path' with the path to your folder containing files
    folder_path = "path/to/your/folder"
    summarize_folder(folder_path)