# RAG pipeline using LangChain, LLM, and MongoDB Atlas (Vector Search)

# Install required libraries:
# pip install langchain langchain-community openai pymongo[tls] pymupdf tiktoken sentence-transformers python-docx python-pptx

import os
from docx import Document as DocxDocument  # Alias to avoid conflicts
from pptx import Presentation  # For PowerPoint presentations
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.schema import Document

# Load environment variables from .env
load_dotenv()

# Pinecone Configuration
pinecone_api_key = os.getenv("PINECONE_API_KEY")
pinecone_environment = os.getenv("PINECONE_ENVIRONMENT")
pinecone_index_name = os.getenv("PINECONE_INDEX_NAME")

# Initialize Pinecone
pinecone_client = Pinecone(api_key=pinecone_api_key)

if pinecone_index_name not in pinecone_client.list_indexes().names():
    pinecone_client.create_index(
        name=pinecone_index_name,
        dimension=384,  # Adjust this based on your embedding model
        metric='cosine',
        spec=ServerlessSpec(
            cloud='aws',
            region=pinecone_environment
        )
    )

index = pinecone_client.Index(pinecone_index_name)

# Initialize ChatGoogleGenerativeAI
chat_model = ChatGoogleGenerativeAI(
    model="gemini-2.0-flash",
    temperature=0,
    max_tokens=None,
    timeout=None,
    max_retries=2,
    api_key=os.getenv("GEMINI_API_KEY")
)

# Initialize HuggingFaceEmbeddings
hf_embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# Load files from a folder (PDFs, Word, PowerPoint, and text files)
def load_files_from_folder(folder_path):
    all_docs = []
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        print(f"Processing file: {filename}")  # Print the document name
        if filename.endswith(".pdf"):
            loader = PyMuPDFLoader(file_path)
            documents = loader.load()
            all_docs.extend([Document(page_content=doc.page_content, metadata={"source": file_path}) for doc in documents])
        elif filename.endswith(".docx"):
            documents = load_word_document(file_path)
            all_docs.extend(documents)
        elif filename.endswith(".pptx"):
            documents = load_powerpoint_presentation(file_path)
            all_docs.extend(documents)
        elif filename.endswith(".txt"):
            documents = load_text_file(file_path)
            all_docs.extend(documents)
    return all_docs

# Load Word documents
def load_word_document(file_path):
    doc = DocxDocument(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
    return [Document(page_content=content, metadata={"source": file_path})]

# Load PowerPoint presentations
def load_powerpoint_presentation(file_path):
    presentation = Presentation(file_path)
    slides_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        slides_content.append("\n".join(slide_text))
    content = "\n".join(slides_content)
    return [Document(page_content=content, metadata={"source": file_path})]

# Load text files
def load_text_file(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()
    return [Document(page_content=content, metadata={"source": file_path})]

# Split documents into chunks using LangChain's RecursiveCharacterTextSplitter
def process_documents(docs, chunk_size=1000, chunk_overlap=200):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_documents(docs)

# Generate embeddings using HuggingFaceEmbeddings
def generate_embedding(text):
    try:
        return hf_embedding_model.embed_query(text)
    except Exception:
        return None

# Store document chunks and embeddings in Pinecone
def store_chunks_in_pinecone(chunks):
    for i, chunk in enumerate(chunks):
        embedding = generate_embedding(chunk.page_content)
        if not embedding:
            continue

        metadata = chunk.metadata if isinstance(chunk.metadata, dict) else {}
        sanitized_metadata = {}
        for key, value in metadata.items():
            if isinstance(value, (str, int, float, bool)):
                sanitized_metadata[key] = value
            elif isinstance(value, list) and all(isinstance(v, str) for v in value):
                sanitized_metadata[key] = value
            else:
                sanitized_metadata[key] = str(value)

        sanitized_metadata["content"] = chunk.page_content

        try:
            index.upsert([(f"doc-{i}", embedding, sanitized_metadata)])
        except Exception:
            pass

# Retrieve relevant documents from Pinecone using Vector Search
def retrieve_documents_from_pinecone(query):
    query_embedding = generate_embedding(query)
    if not query_embedding:
        return []

    results = index.query(
        vector=query_embedding,
        top_k=10,
        include_metadata=True
    )

    return [{"content": match["metadata"]["content"], "metadata": match["metadata"]} for match in results["matches"]]

# Call Gemini API for question-answering using ChatGoogleGenerativeAI
def call_gemini_api(question, context):
    prompt = f"Question: {question}\nContext: {context}\nAnswer:" if context else f"Question: {question}\nAnswer:"
    response = chat_model.invoke(prompt)  # Invoke the Gemini API
    # Extract the 'content' attribute of the response
    return response.content if hasattr(response, "content") else "No answer available."

# Main execution
def main():
    folder_path = "#pa0th"  # Folder containing multiple files
    docs = load_files_from_folder(folder_path)  # Load all supported files
    print(f"Loaded {len(docs)} documents from folder: {folder_path}")

    # Process all documents into chunks
    chunks = process_documents(docs)
    print(f"Created {len(chunks)} chunks from all documents.")

    # Store all chunks and their embeddings in Pinecone
    store_chunks_in_pinecone(chunks)

    while True:
        query = input("Ask a question (or type 'exit'): ")
        if query.lower() == 'exit':
            break

        # Retrieve relevant chunks from Pinecone
        retrieved_docs = retrieve_documents_from_pinecone(query)

        if not retrieved_docs:
            print("No relevant documents found.")
            print("\nAnswer: No relevant documents found.")
            continue

        # Construct context from retrieved chunks
        context = " ".join([doc["content"] for doc in retrieved_docs])

        # Generate an answer using the Gemini API
        answer = call_gemini_api(question=query, context=context)
        print("\nAnswer:", answer)  # Print only the answer content

if __name__ == "__main__":
    main()
